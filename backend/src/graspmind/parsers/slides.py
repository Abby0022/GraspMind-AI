"""PPTX slide parser using python-pptx.

Extracts text from PowerPoint presentations on a per-slide basis.
Each slide becomes a logical "page" in the ParsedDocument.
"""

import logging
from pathlib import Path

from pptx import Presentation

from graspmind.parsers.pdf import ParsedDocument, ParsedPage

logger = logging.getLogger(__name__)


def parse_pptx(file_path: str | Path) -> ParsedDocument:
    """Parse a PPTX file, extracting text per slide.

    Args:
        file_path: Path to the PPTX file.

    Returns:
        ParsedDocument with one page per slide.
    """
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"PPTX not found: {path}")

    try:
        prs = Presentation(str(path))
    except Exception as exc:
        raise ValueError(f"Failed to open PPTX: {exc}") from exc

    pages: list[ParsedPage] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        headings: list[str] = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if not text:
                    continue

                texts.append(text)

                # First text on slide or large text is likely a heading
                if not headings and len(texts) <= 2:
                    headings.append(text)

            # Extract table content
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        texts.append(" | ".join(cells))

        slide_text = "\n".join(texts)
        if slide_text.strip():
            pages.append(ParsedPage(
                page_number=slide_num,
                text=f"[Slide {slide_num}]\n{slide_text}",
                headings=headings,
            ))

    title = path.stem
    if pages and pages[0].headings:
        title = pages[0].headings[0]

    logger.info("Parsed PPTX: %s (%d slides)", title, len(pages))

    return ParsedDocument(
        title=title,
        pages=pages,
        total_pages=len(pages),
        source_type="pptx",
        metadata={
            "file_name": path.name,
            "file_size": path.stat().st_size,
            "slide_count": len(prs.slides),
        },
    )
